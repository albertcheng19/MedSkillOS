#!/usr/bin/env python3
"""
File Security Toolkit
Encrypt/decrypt files, redact privacy info, check password strength.
"""

import argparse
import base64
import io
import os
import re
import sys
from pathlib import Path
from typing import List

import pyzipper
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
from cryptography.fernet import Fernet
from docx import Document
from pptx import Presentation
from PIL import Image


PASSWORD_MIN_LEN = 8
PBKDF2_ITERATIONS = 200_000
FILE_HEADER = b"FSE1"


EMAIL_RE = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b")
PHONE_RE = re.compile(r"\b1[3-9]\d{9}\b")
ID_RE = re.compile(
    r"\b[1-9]\d{5}(19|20)\d{2}(0[1-9]|1[0-2])"
    r"(0[1-9]|[12]\d|3[01])\d{3}[\dXx]\b"
)
KEYWORDS = [
    "Name",
    "address",
    "address",
    "Contact address",
    "Contact person",
    "home address",
]


def check_password_strength(password: str) -> bool:
    if len(password) < PASSWORD_MIN_LEN:
        return False
    has_upper = any(c.isupper() for c in password)
    has_lower = any(c.islower() for c in password)
    has_digit = any(c.isdigit() for c in password)
    return has_upper and has_lower and has_digit


def redact_text(text: str) -> str:
    text = EMAIL_RE.sub("[REDACTED_EMAIL]", text)
    text = PHONE_RE.sub("[REDACTED_PHONE]", text)
    text = ID_RE.sub("[REDACTED_ID]", text)

    for kw in KEYWORDS:
        text = re.sub(rf"({kw}\s*[:：]\s*)([^\n\r]+)", r"\1[REDACTED]", text)
        text = re.sub(rf"({kw}\s+)(\S+)", r"\1[REDACTED]", text)

    return text


def redact_text_file(input_path: Path, output_path: Path) -> None:
    content = input_path.read_text(encoding="utf-8")
    output_path.write_text(redact_text(content), encoding="utf-8")


def redact_docx(input_path: Path, output_path: Path) -> None:
    doc = Document(input_path)
    for p in doc.paragraphs:
        if p.text:
            p.text = redact_text(p.text)
    doc.save(output_path)


def redact_pptx(input_path: Path, output_path: Path) -> None:
    pres = Presentation(input_path)
    for slide in pres.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        run.text = redact_text(run.text)
    pres.save(output_path)


def zip_encrypt(input_path: Path, output_path: Path, password: str) -> None:
    with pyzipper.AESZipFile(
        output_path,
        "w",
        compression=pyzipper.ZIP_DEFLATED,
        encryption=pyzipper.WZ_AES,
    ) as zf:
        zf.setpassword(password.encode("utf-8"))
        if input_path.is_dir():
            for file_path in input_path.rglob("*"):
                if file_path.is_file():
                    arcname = file_path.relative_to(input_path)
                    zf.write(file_path, arcname.as_posix())
        else:
            zf.write(input_path, input_path.name)


def zip_decrypt(input_path: Path, output_path: Path, password: str) -> None:
    output_path.mkdir(parents=True, exist_ok=True)
    with pyzipper.AESZipFile(input_path, "r") as zf:
        zf.setpassword(password.encode("utf-8"))
        zf.extractall(output_path)


def derive_key(password: str, salt: bytes) -> bytes:
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,
        salt=salt,
        iterations=PBKDF2_ITERATIONS,
    )
    return base64.urlsafe_b64encode(kdf.derive(password.encode("utf-8")))


def file_encrypt(input_path: Path, output_path: Path, password: str) -> None:
    data = input_path.read_bytes()
    salt = os.urandom(16)
    key = derive_key(password, salt)
    token = Fernet(key).encrypt(data)
    output_path.write_bytes(FILE_HEADER + salt + token)


def file_decrypt(input_path: Path, output_path: Path, password: str) -> None:
    data = input_path.read_bytes()
    if not data.startswith(FILE_HEADER):
        raise ValueError("Invalid encrypted file header")
    salt = data[len(FILE_HEADER) : len(FILE_HEADER) + 16]
    token = data[len(FILE_HEADER) + 16 :]
    key = derive_key(password, salt)
    raw = Fernet(key).decrypt(token)
    output_path.write_bytes(raw)


def run_redact(input_path: Path, output_path: Path) -> None:
    suffix = input_path.suffix.lower()
    if suffix in {".txt", ".md", ".csv"}:
        redact_text_file(input_path, output_path)
    elif suffix == ".docx":
        redact_docx(input_path, output_path)
    elif suffix == ".pptx":
        redact_pptx(input_path, output_path)
    else:
        raise ValueError("Unsupported file type for redact")


def main() -> None:
    parser = argparse.ArgumentParser(description="File Security Toolkit")
    subparsers = parser.add_subparsers(dest="command", required=True)

    zip_enc = subparsers.add_parser("zip-encrypt")
    zip_enc.add_argument("--input", required=True)
    zip_enc.add_argument("--output", required=True)
    zip_enc.add_argument("--password", required=True)

    zip_dec = subparsers.add_parser("zip-decrypt")
    zip_dec.add_argument("--input", required=True)
    zip_dec.add_argument("--output", required=True)
    zip_dec.add_argument("--password", required=True)

    file_enc = subparsers.add_parser("file-encrypt")
    file_enc.add_argument("--input", required=True)
    file_enc.add_argument("--output", required=True)
    file_enc.add_argument("--password", required=True)

    file_dec = subparsers.add_parser("file-decrypt")
    file_dec.add_argument("--input", required=True)
    file_dec.add_argument("--output", required=True)
    file_dec.add_argument("--password", required=True)

    redact = subparsers.add_parser("redact")
    redact.add_argument("--input", required=True)
    redact.add_argument("--output", required=True)

    check = subparsers.add_parser("check-password")
    check.add_argument("--password", required=True)

    args = parser.parse_args()
    cmd = args.command

    if cmd == "check-password":
        ok = check_password_strength(args.password)
        print("OK" if ok else "WEAK")
        sys.exit(0 if ok else 2)

    if cmd in {"zip-encrypt", "zip-decrypt", "file-encrypt", "file-decrypt"}:
        if not check_password_strength(args.password):
            print("WEAK PASSWORD")
            sys.exit(2)

    input_path = Path(args.input)
    output_path = Path(args.output)

    if cmd == "zip-encrypt":
        zip_encrypt(input_path, output_path, args.password)
    elif cmd == "zip-decrypt":
        zip_decrypt(input_path, output_path, args.password)
    elif cmd == "file-encrypt":
        file_encrypt(input_path, output_path, args.password)
    elif cmd == "file-decrypt":
        file_decrypt(input_path, output_path, args.password)
    elif cmd == "redact":
        run_redact(input_path, output_path)
    else:
        raise ValueError("Unsupported command")

    print("DONE")


if __name__ == "__main__":
    main()
