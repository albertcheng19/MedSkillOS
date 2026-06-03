# -*- coding: utf-8 -*-
"""
Component renderers for converting HTML components to PowerPoint shapes
"""

import io
import urllib.request
import urllib.error
from pathlib import Path

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .layout import LayoutCalculator, px_to_inch, parse_style
from .styles import Colors, FontSizes, Fonts, apply_text_style, apply_shape_fill
from .animation import AnimationBuilder, AnimationResolver


class ComponentRenderers:
    """
    Collection of component renderers.
    Each renderer creates shapes and queues animations for a component type.
    """
    
    def __init__(self, layout: LayoutCalculator, resolver: AnimationResolver, 
                 anim_builder: AnimationBuilder, base_path: Path = None):
        self.layout = layout
        self.resolver = resolver
        self.anim_builder = anim_builder
        # Base path for resolving relative image paths (usually the HTML directory)
        self.base_path = base_path
    
    def render(self, slide, element_data: dict, slide_data: dict):
        """
        Dispatch to the appropriate renderer based on element type.
        
        Args:
            slide: The python-pptx slide object
            element_data: The element configuration from slides-data.js
            slide_data: The parent slide configuration
        """
        component_type = element_data.get('type')
        
        renderers = {
            'comparison': self.render_comparison,
            'terminal': self.render_terminal,
            'quote': self.render_quote,
            'assumptions': self.render_assumptions,
            'timeline': self.render_timeline,
            'stats': self.render_stats,
            'valueCards': self.render_value_cards,
            'competitionBox': self.render_competition_box,
            'strategyBox': self.render_strategy_box,
            'ending': self.render_ending,
            'image': self.render_image,
        }
        
        renderer = renderers.get(component_type)
        if renderer:
            # Parse margins and width in style attribute
            style_str = element_data.get('style', '')
            margins = parse_style(style_str)
            
            # Stores the style information of the current element for use by the component renderer
            self._current_style = margins
            
            # Apply margin-top (increase Y position of fluid layout before rendering)
            if margins.top > 0:
                self.layout.flow.current_y += margins.top
            
            # Record the number of shapes before rendering, used to track newly created shapes
            shapes_before = len(slide.shapes)
            
            # render component
            renderer(slide, element_data, slide_data)
            
            # Clear current style
            self._current_style = None
            
            # Apply margin-bottom (adds extra spacing after rendering)
            if margins.bottom > 0:
                self.layout.flow.current_y += margins.bottom
            
            # Check hideStep to add exit animation for newly created shapes
            hide_step = element_data.get('hideStep')
            if hide_step:
                shapes_after = len(slide.shapes)
                # Add exit animation to all newly created shapes
                for i in range(shapes_before, shapes_after):
                    shape = slide.shapes[i]
                    self.anim_builder.queue_exit(shape, 'fadeOut', hide_step, 0, 300)
        else:
            print(f"Warning: Unknown component type '{component_type}'")
    
    # ===== Component Renderers =====
    
    def render_comparison(self, slide, element_data: dict, slide_data: dict):
        """Render comparison cards (left/right)"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        stagger = anim_config.get('stagger', 0) or 30
        
        # Get custom width (if any)
        custom_width = getattr(self, '_current_style', None)
        width = custom_width.width if custom_width and custom_width.width else None
        
        # Get the number of left/right items for dynamic calculation of height
        left_items = len(element_data.get('left', {}).get('items', []))
        right_items = len(element_data.get('right', {}).get('items', []))
        
        left_rect, right_rect, height = self.layout.comparison_cards(
            custom_width=width, left_items=left_items, right_items=right_items
        )
        self.layout.flow.advance(height)
        
        # Left card
        left_shapes = self._create_comparison_card(
            slide, element_data.get('left', {}), left_rect
        )
        for i, shape in enumerate(left_shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * stagger, duration)
        
        # Right card
        right_shapes = self._create_comparison_card(
            slide, element_data.get('right', {}), right_rect
        )
        base_delay = len(left_shapes) * stagger
        for i, shape in enumerate(right_shapes):
            self.anim_builder.queue_entrance(shape, preset, step, 
                                             base_delay + i * stagger, duration)
    
    def _create_comparison_card(self, slide, card_data: dict, rect) -> list:
        """Create a single comparison card and return all shapes
        CSS: .comparison-card { padding: 16px 18px; }
        CSS: .comparison-card h4 { font-size: 16px; margin-bottom: 10px; padding-bottom: 8px; border-bottom: 1px solid var(--border-color); }
        CSS: .comparison-list li { padding: 6px 0; font-size: 14px; }
        CSS: .comparison-list li::before { content: '→'; color: var(--accent-coral); }
        
        Fix: Increase item height to support line wrapping of long text"""
        shapes = []
        
        # Card background - with shadow effect
        x, y, w, h = rect.to_inches()
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        apply_shape_fill(card, Colors.BG_CARD, Colors.BORDER)
        shapes.append(card)
        
        # Internal margin - CSS: padding: 16px 18px
        pad_x = Inches(px_to_inch(18))
        pad_y = Inches(px_to_inch(16))
        
        # Card title with icon - CSS: font-size: 16px
        title_height = Inches(px_to_inch(22))
        icon = card_data.get('icon', '')
        title_text = card_data.get('title', '')
        
        # Icons are rendered individually (using accent colors)
        icon_width = Inches(px_to_inch(20))
        if icon:
            icon_box = slide.shapes.add_textbox(
                x + pad_x, y + pad_y,
                icon_width, title_height
            )
            apply_text_style(
                icon_box.text_frame, icon,
                font_size=FontSizes.CARD_TITLE,
                font_color=Colors.ACCENT_CORAL,
                bold=True
            )
            shapes.append(icon_box)
        
        # title text
        title_box = slide.shapes.add_textbox(
            x + pad_x + (icon_width if icon else Inches(0)), y + pad_y,
            w - pad_x * 2 - (icon_width if icon else Inches(0)), title_height
        )
        apply_text_style(
            title_box.text_frame, title_text,
            font_size=FontSizes.CARD_TITLE,
            font_color=Colors.TEXT_PRIMARY,
            bold=True,
            word_wrap=True
        )
        shapes.append(title_box)
        
        # Separator line below title - CSS: border-bottom: 1px solid var(--border-color)
        divider_y = y + pad_y + title_height + Inches(px_to_inch(8))
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + pad_x, divider_y,
            w - pad_x * 2, Inches(px_to_inch(1))
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = Colors.BORDER
        divider.line.fill.background()
        shapes.append(divider)
        
        # Card items - CSS: padding: 6px 0; font-size: 14px
        # Fix: Increase item height to support long text wrapping (from 18px to 36px)
        items = card_data.get('items', [])
        item_start_y = Inches(px_to_inch(50))  # title + divider + gap
        item_height = Inches(px_to_inch(36))  # Increase height to support 2 lines of text
        item_gap = Inches(px_to_inch(6))
        arrow_width = Inches(px_to_inch(16))
        
        for i, item in enumerate(items):
            item_y = y + item_start_y + (item_height + item_gap) * i
            
            # Arrows (use accent color) - CSS: .comparison-list li::before { color: var(--accent-coral); }
            arrow_box = slide.shapes.add_textbox(
                x + pad_x, item_y,
                arrow_width, item_height
            )
            apply_text_style(
                arrow_box.text_frame, "→",
                font_size=FontSizes.CARD_ITEM,
                font_color=Colors.ACCENT_CORAL,
                bold=True
            )
            shapes.append(arrow_box)
            
            # Text content - enable word wrapping
            item_box = slide.shapes.add_textbox(
                x + pad_x + arrow_width, item_y,
                w - pad_x * 2 - arrow_width, item_height
            )
            apply_text_style(
                item_box.text_frame, item,
                font_size=FontSizes.CARD_ITEM,
                font_color=Colors.TEXT_SECONDARY,
                word_wrap=True  # Enable word wrapping
            )
            shapes.append(item_box)
        
        return shapes
    
    def render_terminal(self, slide, element_data: dict, slide_data: dict):
        """Render terminal/code block
        CSS: .terminal-header { padding: 12px 16px; background: #2d2d2d; }
        CSS: .terminal-dot { width: 12px; height: 12px; border-radius: 50%; }
        CSS: .terminal-dot.red { background: #ff5f56; }
        CSS: .terminal-dot.yellow { background: #ffbd2e; }
        CSS: .terminal-dot.green { background: #27c93f; }
        CSS: .terminal-body { padding: 16px 20px; font-size: 13px; line-height: 2; }
        
        Fix: Add 3 colored dots, dynamically calculate height, enable automatic word wrapping"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        # Count the number of lines of content (count newlines and estimate text length)
        content_parts = element_data.get('content', [])
        full_text = ''.join(part.get('text', '') for part in content_parts if part.get('type') != 'cursor')
        newline_count = full_text.count('\n')
        # Estimated text wrapping: every 50 characters
        char_lines = len(full_text.replace('\n', '')) // 50
        content_lines = max(3, newline_count + char_lines + 2)  # At least 3 lines
        
        # Get custom width (if any)
        custom_width = getattr(self, '_current_style', None)
        width = custom_width.width if custom_width and custom_width.width else None
        
        rect, height = self.layout.terminal(custom_width=width, content_lines=content_lines)
        self.layout.flow.advance(height)
        shapes = []
        
        # Terminal background - CSS: background: #3d3d3d; border-radius: 16px
        x, y, w, h = rect.to_inches()
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        apply_shape_fill(bg, Colors.TERMINAL_BG, Colors.TERMINAL_BG)
        shapes.append(bg)
        
        # Terminal header - CSS: padding: 12px 16px; background: #2d2d2d
        header_height = Inches(px_to_inch(36))  # 12px padding * 2 + 12px dot
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, w, header_height
        )
        apply_shape_fill(header, Colors.TERMINAL_HEADER, Colors.TERMINAL_HEADER)
        shapes.append(header)
        
        # Terminal dots (red, yellow, green) - CSS: .terminal-dot { width: 12px; height: 12px; }
        dot_size = Inches(px_to_inch(12))
        dot_gap = Inches(px_to_inch(8))  # gap: 8px
        dot_start_x = x + Inches(px_to_inch(16))  # padding-left: 16px
        dot_y = y + Inches(px_to_inch(12))  # padding-top: 12px
        
        # Red dot - #ff5f56
        red_dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            dot_start_x, dot_y, dot_size, dot_size
        )
        red_dot.fill.solid()
        red_dot.fill.fore_color.rgb = RGBColor(0xFF, 0x5F, 0x56)
        red_dot.line.fill.background()
        shapes.append(red_dot)
        
        # Yellow dot - #ffbd2e
        yellow_dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            dot_start_x + dot_size + dot_gap, dot_y, dot_size, dot_size
        )
        yellow_dot.fill.solid()
        yellow_dot.fill.fore_color.rgb = RGBColor(0xFF, 0xBD, 0x2E)
        yellow_dot.line.fill.background()
        shapes.append(yellow_dot)
        
        # Green dot - #27c93f
        green_dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            dot_start_x + (dot_size + dot_gap) * 2, dot_y, dot_size, dot_size
        )
        green_dot.fill.solid()
        green_dot.fill.fore_color.rgb = RGBColor(0x27, 0xC9, 0x3F)
        green_dot.line.fill.background()
        shapes.append(green_dot)
        
        # Content text - CSS: padding: 16px 20px; font-size: 13px; line-height: 2
        pad_x = Inches(px_to_inch(20))  # padding: 16px 20px
        pad_y = Inches(px_to_inch(16))
        text_box = slide.shapes.add_textbox(
            x + pad_x, y + header_height + pad_y,
            w - pad_x * 2, h - header_height - pad_y * 2
        )
        
        # Enable word wrapping
        text_box.text_frame.word_wrap = True
        
        # Build content with multiple runs for different colors
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.font.name = Fonts.MONO
        p.font.size = FontSizes.TERMINAL_TEXT
        
        # Process each part with appropriate color
        first_run = True
        for part in content_parts:
            part_type = part.get('type', 'text')
            part_text = part.get('text', '')
            
            if part_type == 'cursor':
                continue  # Skip cursor markers
            
            if first_run:
                # Use the first paragraph's default run
                run = p.runs[0] if p.runs else p.add_run()
                run.text = part_text
                first_run = False
            else:
                run = p.add_run()
                run.text = part_text
            
            # Apply color based on type
            run.font.name = Fonts.MONO
            run.font.size = FontSizes.TERMINAL_TEXT
            
            if part_type == 'prompt':
                run.font.color.rgb = Colors.TERMINAL_PROMPT  # #D97757
            elif part_type == 'highlight':
                run.font.color.rgb = Colors.TERMINAL_HIGHLIGHT  # #D97757
            elif part_type == 'success':
                run.font.color.rgb = Colors.ACCENT_GREEN  # #5a9a6a
            else:
                run.font.color.rgb = Colors.TERMINAL_TEXT  # #e0e0e0
        
        shapes.append(text_box)
        
        # Queue animations
        for i, shape in enumerate(shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * 20, duration)
    
    def render_quote(self, slide, element_data: dict, slide_data: dict):
        """Render quote block
        CSS: .quote-block { padding: 20px 28px; background: gradient; }
        CSS: .quote-block::before { content: '"'; position: absolute; font-size: 36px; color: var(--accent-coral); opacity: 0.25; }
        CSS: .quote-text { font-size: 18px; font-weight: 600; font-style: italic; }
        CSS: .quote-author { margin-top: 12px; font-size: 13px; }
        
        Fix: Support multi-line quoted text, enable automatic word wrapping"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        # Check whether newline characters are included and decide whether to use a larger height
        text = element_data.get('text', '')
        is_multiline = '\n' in text or len(text) > 60  # Contains line breaks or long text
        
        rect, height = self.layout.quote(multiline=is_multiline)
        self.layout.flow.advance(height)
        shapes = []
        
        # Quote background - CSS: background: linear-gradient
        x, y, w, h = rect.to_inches()
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        # Use a light orange background to approximate a gradient effect
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFD, 0xF8, 0xF5)  # very light orange
        bg.line.color.rgb = Colors.BORDER
        bg.line.width = Pt(1)
        shapes.append(bg)
        
        # Internal margin - CSS: padding: 20px 28px
        pad_x = Inches(px_to_inch(28))
        pad_y = Inches(px_to_inch(20))
        
        # Quote decoration (") - CSS: .quote-block::before
        quote_mark = slide.shapes.add_textbox(
            x + Inches(px_to_inch(14)), y + Inches(px_to_inch(8)),
            Inches(px_to_inch(30)), Inches(px_to_inch(36))
        )
        quote_mark.text_frame.paragraphs[0].text = '"'
        quote_mark.text_frame.paragraphs[0].font.size = Pt(24)
        quote_mark.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xE8, 0xC0, 0xAA)  # light orange (opacity: 0.25 effect)
        quote_mark.text_frame.paragraphs[0].font.name = Fonts.DISPLAY
        shapes.append(quote_mark)
        
        # Quote text - CSS: font-size: 18px
        # CSS: .quote-text { font-style: italic; }
        # CSS: .quote-text.gradient { font-style: normal; }
        is_gradient = element_data.get('gradient', False)
        text_color = Colors.ACCENT_CORAL if is_gradient else Colors.TEXT_PRIMARY
        # When gradient=true, don't use italic (CSS: .quote-text.gradient { font-style: normal; })
        use_italic = not is_gradient
        
        author = element_data.get('author')
        author_height = Inches(px_to_inch(22)) if author else Inches(0)
        
        text_box = slide.shapes.add_textbox(
            x + pad_x, y + pad_y,
            w - pad_x * 2, h - pad_y * 2 - author_height
        )
        apply_text_style(
            text_box.text_frame, text,
            font_size=FontSizes.QUOTE_TEXT,
            font_color=text_color,
            bold=True, italic=use_italic,
            alignment=PP_ALIGN.CENTER,
            word_wrap=True  # Enable word wrapping
        )
        shapes.append(text_box)
        
        # Author (if present) - CSS: font-size: 13px
        if author:
            author_icon = author.get('icon', '')
            author_text_str = author.get('text', '')
            # Use accent color for icons
            if author_icon:
                author_full = f"{author_icon} {author_text_str}"
            else:
                author_full = author_text_str
            
            author_box = slide.shapes.add_textbox(
                x + pad_x, y + h - pad_y - author_height,
                w - pad_x * 2, author_height
            )
            apply_text_style(
                author_box.text_frame, author_full,
                font_size=FontSizes.QUOTE_AUTHOR,
                font_color=Colors.TEXT_TERTIARY,
                alignment=PP_ALIGN.RIGHT,
                word_wrap=True  # Enable word wrapping
            )
            shapes.append(author_box)
        
        # Queue animations
        for i, shape in enumerate(shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * 40, duration)
    
    def render_assumptions(self, slide, element_data: dict, slide_data: dict):
        """Render assumption grid with two-phase reveal
        CSS: .assumption-item { padding: 14px 18px; }
        CSS: .assumption-old { font-size: 14px; margin-bottom: 8px; }
        CSS: .assumption-new { font-size: 14px; font-weight: 500; }
        
        Fix: Increase new_text height to support long text wrapping"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        reveal_step = element_data.get('revealStep', step + 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        stagger = anim_config.get('stagger', 0) or 100
        
        items = element_data.get('items', [])
        rects, height = self.layout.assumptions_grid(len(items))
        self.layout.flow.advance(height)
        
        # CSS Dimensions - Increase new_text height to support wrapping of long text
        pad_x = Inches(px_to_inch(18))  # padding: 14px 18px
        pad_y = Inches(px_to_inch(14))
        old_text_height = Inches(px_to_inch(22))
        new_text_height = Inches(px_to_inch(55))  # Increased height to support 2-3 lines of text (from 22px to 55px)
        gap = Inches(px_to_inch(8))
        
        for i, (item, rect) in enumerate(zip(items, rects)):
            x, y, w, h = rect.to_inches()
            base_delay = i * stagger
            
            # Card background - CSS: border: 1px solid var(--border-color)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            apply_shape_fill(card, Colors.BG_CARD, Colors.BORDER)
            self.anim_builder.queue_entrance(card, preset, step, base_delay, duration)
            
            # Old text (shows at step) - CSS: font-size: 14px; color: var(--text-tertiary)
            old_text = item.get('old', '')
            old_box = slide.shapes.add_textbox(
                x + pad_x, y + pad_y,
                w - pad_x * 2, old_text_height
            )
            apply_text_style(
                old_box.text_frame, old_text,
                font_size=FontSizes.ASSUMPTION_OLD,
                font_color=Colors.TEXT_TERTIARY,
                word_wrap=True
            )
            self.anim_builder.queue_entrance(old_box, preset, step, 
                                             base_delay + 20, duration)
            
            # New text (shows at revealStep) - CSS: font-size: 14px; color: var(--accent-green)
            # Increase text box height to support line wrapping of long text
            new_text = item.get('new', '')
            new_box = slide.shapes.add_textbox(
                x + pad_x, y + pad_y + old_text_height + gap,
                w - pad_x * 2, new_text_height
            )
            apply_text_style(
                new_box.text_frame, new_text,
                font_size=FontSizes.ASSUMPTION_NEW,
                font_color=Colors.ACCENT_GREEN,
                bold=True,
                word_wrap=True  # Enable word wrapping
            )
            # New text appears at revealStep
            self.anim_builder.queue_entrance(new_box, 'slideLeft', reveal_step,
                                             i * 80, duration)
    
    def render_timeline(self, slide, element_data: dict, slide_data: dict):
        """Render timeline with progressive activation
        CSS: .timeline-item { padding: 14px; }
        CSS: .timeline-badge { padding: 4px 12px; background: rgba(217, 119, 87, 0.1); border-radius: 100px; }
        CSS: .timeline-badge.current { background: rgba(90, 154, 106, 0.1); color: var(--accent-green); }
        CSS: .timeline-icon { font-size: 24px; margin-bottom: 8px; }
        CSS: .timeline-title { font-size: 14px; margin-bottom: 4px; }
        CSS: .timeline-desc { font-size: 12px; }
        CSS: .timeline-case { font-size: 11px; margin-top: 6px; font-style: italic; }
        
        Fix: Increase case height to support line wrapping of long text"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        stagger = anim_config.get('stagger', 0) or 100
        
        items = element_data.get('items', [])
        rects, height = self.layout.timeline_items(len(items))
        self.layout.flow.advance(height)
        
        # CSS Dimensions - Increase case height to support wrapping of long text
        pad = Inches(px_to_inch(14))
        badge_height = Inches(px_to_inch(22))  # Contains padding
        icon_height = Inches(px_to_inch(28))
        title_height = Inches(px_to_inch(18))
        desc_height = Inches(px_to_inch(16))
        case_height = Inches(px_to_inch(38))  # Increased height to support 2 lines of text (from 24px to 38px)
        gap_small = Inches(px_to_inch(4))
        gap_medium = Inches(px_to_inch(8))
        
        for i, (item, rect) in enumerate(zip(items, rects)):
            x, y, w, h = rect.to_inches()
            active_step = item.get('activeStep', step)
            is_badge_current = item.get('badgeCurrent', False)
            is_icon_filled = item.get('iconFilled', False)
            # Card border highlight when either badgeCurrent or iconFilled
            has_highlight = is_badge_current or is_icon_filled
            
            base_delay = i * stagger
            
            # Item background - CSS: border: 1px solid var(--border-color)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            if has_highlight:
                apply_shape_fill(card, Colors.BG_CARD, Colors.ACCENT_CORAL, Pt(1.5))
            else:
                apply_shape_fill(card, Colors.BG_CARD, Colors.BORDER)
            
            current_y = y + pad
            
            # Badge background (pill shape) - CSS: background: rgba(217, 119, 87, 0.1); border-radius: 100px
            badge_text = item.get('badge', '')
            badge_width = Inches(px_to_inch(60))
            badge_x = x + (w - badge_width) / 2
            
            # Background pill shape
            badge_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                badge_x, current_y,
                badge_width, badge_height
            )
            # Light background - orange or green
            if is_badge_current:
                badge_bg.fill.solid()
                badge_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)  # light green
            else:
                badge_bg.fill.solid()
                badge_bg.fill.fore_color.rgb = RGBColor(0xFE, 0xF0, 0xEB)  # light orange
            badge_bg.line.fill.background()
            
            # Badge text
            badge_box = slide.shapes.add_textbox(
                badge_x, current_y,
                badge_width, badge_height
            )
            badge_color = Colors.ACCENT_GREEN if is_badge_current else Colors.ACCENT_CORAL
            apply_text_style(
                badge_box.text_frame, badge_text,
                font_size=FontSizes.TIMELINE_BADGE,
                font_color=badge_color,
                bold=True,
                alignment=PP_ALIGN.CENTER
            )
            current_y += badge_height + gap_medium
            
            # Icon - CSS: font-size: 24px
            # CSS: .timeline-icon.claude-icon.filled { color: var(--accent-green); }
            icon_text = item.get('icon', '')
            icon_box = slide.shapes.add_textbox(
                x + pad, current_y,
                w - pad * 2, icon_height
            )
            # iconFilled=true -> green, otherwise coral if highlighted, else tertiary
            if is_icon_filled:
                icon_color = Colors.ACCENT_GREEN
            elif has_highlight:
                icon_color = Colors.ACCENT_CORAL
            else:
                icon_color = Colors.TEXT_TERTIARY
            apply_text_style(
                icon_box.text_frame, icon_text,
                font_size=FontSizes.TIMELINE_ICON,
                font_color=icon_color,
                alignment=PP_ALIGN.CENTER
            )
            current_y += icon_height + gap_medium
            
            # Title - CSS: font-size: 14px; font-weight: 600
            title_text = item.get('title', '')
            title_box = slide.shapes.add_textbox(
                x + pad, current_y,
                w - pad * 2, title_height
            )
            apply_text_style(
                title_box.text_frame, title_text,
                font_size=FontSizes.TIMELINE_TITLE,
                font_color=Colors.TEXT_PRIMARY,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                word_wrap=True
            )
            current_y += title_height + gap_small
            
            # Description - CSS: font-size: 12px
            desc_text = item.get('desc', '')
            desc_box = slide.shapes.add_textbox(
                x + pad, current_y,
                w - pad * 2, desc_height
            )
            apply_text_style(
                desc_box.text_frame, desc_text,
                font_size=FontSizes.TIMELINE_DESC,
                font_color=Colors.TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
                word_wrap=True
            )
            current_y += desc_height + gap_small
            
            # Case (if present) - CSS: font-size: 11px; font-style: italic
            # Fix: Enable word wrapping to support long text
            case_text = item.get('case', '')
            case_box = None
            if case_text:
                case_box = slide.shapes.add_textbox(
                    x + pad, current_y,
                    w - pad * 2, case_height
                )
                apply_text_style(
                    case_box.text_frame, case_text,
                    font_size=FontSizes.TIMELINE_CASE,
                    font_color=Colors.TEXT_TERTIARY,
                    italic=True,
                    alignment=PP_ALIGN.CENTER,
                    word_wrap=True  # Enable word wrapping
                )
            
            # Queue all shapes for this item's activeStep
            all_shapes = [card, badge_bg, badge_box, icon_box, title_box, desc_box]
            if case_box:
                all_shapes.append(case_box)
            for j, shape in enumerate(all_shapes):
                self.anim_builder.queue_entrance(shape, preset, active_step,
                                                 base_delay + j * 15, duration)
    
    def render_stats(self, slide, element_data: dict, slide_data: dict):
        """Render stats row
        CSS: .stats-title { font-size: 14px; margin-bottom: 16px; }
        CSS: .stat-item { padding: 16px 20px; }
        CSS: .stat-icon-wrapper { width: 40px; height: 40px; background: gradient; border-radius: 50%; }
        CSS: .stat-value { font-size: 22px; font-weight: 800; gradient text; }
        CSS: .stat-label { font-size: 11px; }
        
        Fix: Increase label height to support line wrapping of long text"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        stagger = anim_config.get('stagger', 0) or 150
        
        items = element_data.get('items', [])
        title = element_data.get('title', '')
        
        # Render the title first (if there is one)
        title_height = Inches(px_to_inch(24)) if title else Inches(0)
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0), Inches(self.layout.flow.get_y()),
                Inches(self.layout.slide_width), title_height
            )
            apply_text_style(
                title_box.text_frame, title,
                font_size=FontSizes.STATS_TITLE,
                font_color=Colors.TEXT_SECONDARY,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                word_wrap=True
            )
            self.anim_builder.queue_entrance(title_box, preset, step, 0, duration)
            self.layout.flow.current_y += px_to_inch(30)
        
        rects, height = self.layout.stats_items(len(items))
        self.layout.flow.advance(height)
        
        # CSS dimensions - increase label height to support wrapping of long text
        pad = Inches(px_to_inch(16))
        icon_wrapper_size = Inches(px_to_inch(40))  # width: 40px; height: 40px
        value_height = Inches(px_to_inch(28))
        label_height = Inches(px_to_inch(36))  # Increased height to support 2 lines of text (from 24px to 36px)
        gap = Inches(px_to_inch(10))
        
        for i, (item, rect) in enumerate(zip(items, rects)):
            x, y, w, h = rect.to_inches()
            base_delay = i * stagger + (100 if title else 0)
            
            # Card background
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            apply_shape_fill(card, Colors.BG_CARD, Colors.BORDER)
            self.anim_builder.queue_entrance(card, preset, step, base_delay, duration)
            
            current_y = y + pad
            icon_center_x = x + (w - icon_wrapper_size) / 2
            
            # Icon wrapper (circle background) - CSS: background: gradient; border-radius: 50%
            icon_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                icon_center_x, current_y,
                icon_wrapper_size, icon_wrapper_size
            )
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = RGBColor(0xFE, 0xF0, 0xEB)  # light orange background
            icon_bg.line.fill.background()
            self.anim_builder.queue_entrance(icon_bg, preset, step, base_delay + 15, duration)
            
            # Icon text - CSS: .stat-icon { font-size: 18px; }
            icon_text = item.get('icon', '')
            icon_box = slide.shapes.add_textbox(
                icon_center_x, current_y,
                icon_wrapper_size, icon_wrapper_size
            )
            apply_text_style(
                icon_box.text_frame, icon_text,
                font_size=FontSizes.STATS_ICON,
                font_color=Colors.ACCENT_CORAL,
                alignment=PP_ALIGN.CENTER
            )
            self.anim_builder.queue_entrance(icon_box, preset, step, base_delay + 20, duration)
            current_y += icon_wrapper_size + gap
            
            # Value - CSS: font-size: 22px; font-weight: 800; gradient text
            value_text = item.get('value', '')
            value_box = slide.shapes.add_textbox(x + pad, current_y, w - pad * 2, value_height)
            apply_text_style(
                value_box.text_frame, value_text,
                font_size=FontSizes.STATS_VALUE,
                font_color=Colors.ACCENT_CORAL,
                bold=True,
                alignment=PP_ALIGN.CENTER
            )
            self.anim_builder.queue_entrance(value_box, preset, step, base_delay + 40, duration)
            current_y += value_height + Inches(px_to_inch(4))
            
            # Label - CSS: font-size: 11px
            # Fix: Enable word wrapping to support long label text
            label_text = item.get('label', '')
            label_box = slide.shapes.add_textbox(x + pad, current_y, w - pad * 2, label_height)
            apply_text_style(
                label_box.text_frame, label_text,
                font_size=FontSizes.STATS_LABEL,
                font_color=Colors.TEXT_TERTIARY,
                alignment=PP_ALIGN.CENTER,
                word_wrap=True  # Enable word wrapping
            )
            self.anim_builder.queue_entrance(label_box, preset, step, base_delay + 60, duration)
    
    def render_value_cards(self, slide, element_data: dict, slide_data: dict):
        """Render value cards with progressive activation
        CSS: .value-card { padding: 16px 12px; border: 1px solid var(--border-color); }
        CSS: .value-card.active { border-color: var(--accent-coral); box-shadow: var(--shadow-glow-strong); }
        CSS: .value-card-icon { width: 48px; height: 48px; background: gradient; border-radius: 50%; }
        CSS: .value-card h4 { font-size: 15px; margin-bottom: 6px; }
        CSS: .value-card p { font-size: 13px; }
        
        Fix: Increase desc height to support line wrapping of long text"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        stagger = anim_config.get('stagger', 0) or 150
        
        items = element_data.get('items', [])
        rects, height = self.layout.value_cards(len(items))
        self.layout.flow.advance(height)
        
        # CSS Dimensions - Increase desc height to support wrapping of long text
        pad_x = Inches(px_to_inch(12))
        pad_y = Inches(px_to_inch(16))
        icon_size = Inches(px_to_inch(44))  # width: 48px (slightly smaller)
        title_h = Inches(px_to_inch(20))
        desc_h = Inches(px_to_inch(48))  # Increased height to support 2 lines of text (from 32px to 48px)
        gap = Inches(px_to_inch(8))
        
        for i, (item, rect) in enumerate(zip(items, rects)):
            x, y, w, h = rect.to_inches()
            active_step = item.get('activeStep', step)
            base_delay = i * stagger
            
            # Card background - CSS: border: 1px solid var(--accent-coral) for active
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            apply_shape_fill(card, Colors.BG_CARD, Colors.ACCENT_CORAL, Pt(1))
            
            current_y = y + pad_y
            icon_center_x = x + (w - icon_size) / 2
            
            # Icon background (circle) - CSS: background: gradient; border-radius: 50%
            icon_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                icon_center_x, current_y,
                icon_size, icon_size
            )
            icon_bg.fill.solid()
            icon_bg.fill.fore_color.rgb = RGBColor(0xFE, 0xF0, 0xEB)  # light orange
            icon_bg.line.fill.background()
            
            # Icon text - CSS: font-size: 20px
            icon_text = item.get('icon', '')
            icon_box = slide.shapes.add_textbox(
                icon_center_x, current_y,
                icon_size, icon_size
            )
            apply_text_style(
                icon_box.text_frame, icon_text,
                font_size=FontSizes.VALUE_CARD_ICON,
                font_color=Colors.ACCENT_CORAL,
                alignment=PP_ALIGN.CENTER
            )
            current_y += icon_size + gap
            
            # Title - CSS: font-size: 15px; font-weight: 600
            title_text = item.get('title', '')
            title_box = slide.shapes.add_textbox(
                x + pad_x, current_y,
                w - pad_x * 2, title_h
            )
            apply_text_style(
                title_box.text_frame, title_text,
                font_size=FontSizes.VALUE_CARD_TITLE,
                font_color=Colors.TEXT_PRIMARY,
                bold=True,
                alignment=PP_ALIGN.CENTER
            )
            current_y += title_h + gap
            
            # Description - CSS: font-size: 13px
            # Fix: Enable word wrapping to support long description text
            desc_text = item.get('desc', '')
            desc_box = slide.shapes.add_textbox(
                x + pad_x, current_y,
                w - pad_x * 2, desc_h
            )
            apply_text_style(
                desc_box.text_frame, desc_text,
                font_size=FontSizes.VALUE_CARD_DESC,
                font_color=Colors.TEXT_SECONDARY,
                alignment=PP_ALIGN.CENTER,
                word_wrap=True  # Enable word wrapping
            )
            
            # Queue animations at activeStep
            for j, shape in enumerate([card, icon_bg, icon_box, title_box, desc_box]):
                self.anim_builder.queue_entrance(shape, preset, active_step,
                                                 base_delay + j * 25, duration)
    
    def render_competition_box(self, slide, element_data: dict, slide_data: dict):
        """Render competition analysis box
        CSS: .competition-box { padding: 16px 26px; }
        CSS: .competition-box h4 { font-size: 14px; margin-bottom: 10px; }
        CSS: .competition-item-box { padding: 6px 14px; font-size: 13px; }
        CSS: .competition-conclusion { font-size: 16px; font-weight: 600; }
        
        Fix: Increase items height and enable autoscaling to prevent long text from overflowing"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        rect, height = self.layout.competition_box()
        self.layout.flow.advance(height)
        x, y, w, h = rect.to_inches()
        shapes = []
        
        # CSS dimensions - increase items height to support long text
        pad_x = Inches(px_to_inch(26))  # padding: 16px 26px
        pad_y = Inches(px_to_inch(16))
        title_height = Inches(px_to_inch(20))
        items_height = Inches(px_to_inch(42))  # Increased height to support 2 lines of text (from 28px to 42px)
        conclusion_height = Inches(px_to_inch(22))
        gap = Inches(px_to_inch(10))
        
        # Background - CSS: gradient background
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        # Approximate gradient background
        apply_shape_fill(bg, RGBColor(0xFE, 0xF8, 0xF5), Colors.BORDER)
        shapes.append(bg)
        
        current_y = y + pad_y
        
        # Title - CSS: font-size: 14px
        title_text = element_data.get('title', '')
        title_box = slide.shapes.add_textbox(
            x + pad_x, current_y,
            w - pad_x * 2, title_height
        )
        apply_text_style(
            title_box.text_frame, title_text,
            font_size=FontSizes.COMPETITION_TITLE,
            font_color=Colors.TEXT_PRIMARY,
            alignment=PP_ALIGN.CENTER,
            word_wrap=True
        )
        shapes.append(title_box)
        current_y += title_height + gap
        
        # Items (horizontal) - CSS: font-size: 13px
        # Fix: Enable word wrapping and auto scaling to prevent overflow
        items = element_data.get('items', [])
        items_text = '  |  '.join(items)
        items_box = slide.shapes.add_textbox(
            x + pad_x, current_y,
            w - pad_x * 2, items_height
        )
        apply_text_style(
            items_box.text_frame, items_text,
            font_size=FontSizes.COMPETITION_ITEM,
            font_color=Colors.TEXT_SECONDARY,
            alignment=PP_ALIGN.CENTER,
            word_wrap=True,  # Enable word wrapping
            auto_size='fit'  # Enable automatic scaling
        )
        shapes.append(items_box)
        current_y += items_height + gap
        
        # Conclusion - CSS: font-size: 16px; font-weight: 600
        conclusion_text = element_data.get('conclusion', '')
        conclusion_box = slide.shapes.add_textbox(
            x + pad_x, current_y,
            w - pad_x * 2, conclusion_height
        )
        apply_text_style(
            conclusion_box.text_frame, conclusion_text,
            font_size=FontSizes.COMPETITION_CONCLUSION,
            font_color=Colors.ACCENT_CORAL,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            word_wrap=True
        )
        shapes.append(conclusion_box)
        
        # Queue animations
        for i, shape in enumerate(shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * 50, duration)
    
    def render_strategy_box(self, slide, element_data: dict, slide_data: dict):
        """Render strategy box with right/wrong items
        CSS: .strategy-box { padding: 18px 26px; }
        CSS: .strategy-box h3 { font-size: 16px; margin-bottom: 14px; }
        CSS: .strategy-item { padding: 10px 14px; margin-bottom: 8px; font-size: 14px; }
        CSS: .strategy-final p { font-size: 16px; }
        
        Fix: Increase item height and enable word wrapping to support long text"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        items = element_data.get('items', [])
        rect, height = self.layout.strategy_box(items_count=len(items))
        self.layout.flow.advance(height)
        x, y, w, h = rect.to_inches()
        shapes = []
        
        # CSS dimensions - increase item height to support wrapping of long text
        pad_x = Inches(px_to_inch(26))  # padding: 18px 26px
        pad_y = Inches(px_to_inch(18))
        title_height = Inches(px_to_inch(22))
        item_height = Inches(px_to_inch(36))  # Increased height to support 2 lines of text (from 28px to 36px)
        item_gap = Inches(px_to_inch(8))
        icon_width = Inches(px_to_inch(24))
        final_height = Inches(px_to_inch(24))
        
        # Background - CSS: linear-gradient(135deg, #1a1a1a 0%, #2d2d2d 100%)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        apply_shape_fill(bg, Colors.DARK_BG, Colors.DARK_BG)
        shapes.append(bg)
        
        current_y = y + pad_y
        
        # Title - CSS: font-size: 16px
        title_text = element_data.get('title', '')
        title_box = slide.shapes.add_textbox(
            x + pad_x, current_y,
            w - pad_x * 2, title_height
        )
        apply_text_style(
            title_box.text_frame, title_text,
            font_size=FontSizes.STRATEGY_TITLE,
            font_color=Colors.TEXT_LIGHT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            word_wrap=True
        )
        shapes.append(title_box)
        current_y += title_height + Inches(px_to_inch(14))
        
        # Items - CSS: font-size: 14px
        # CSS: .strategy-item.wrong { opacity: 0.5; }
        # Fix: Enable word wrapping to support long text
        for i, item in enumerate(items):
            item_type = item.get('type', 'wrong')
            item_text = item.get('text', '')
            
            icon = '✗' if item_type == 'wrong' else '✓'
            
            # Use dimmed colors for 'wrong' items to simulate opacity: 0.5
            if item_type == 'wrong':
                icon_color = Colors.CORAL_DIMMED  # Dimmed coral for wrong icon
                text_color = Colors.WHITE_DIMMED  # Dimmed white for wrong text
            else:
                icon_color = Colors.ACCENT_GREEN
                text_color = Colors.WHITE
            
            icon_box = slide.shapes.add_textbox(
                x + pad_x, current_y,
                icon_width, item_height
            )
            apply_text_style(
                icon_box.text_frame, icon,
                font_size=FontSizes.STRATEGY_ITEM,
                font_color=icon_color
            )
            shapes.append(icon_box)
            
            text_box = slide.shapes.add_textbox(
                x + pad_x + icon_width, current_y,
                w - pad_x * 2 - icon_width, item_height
            )
            apply_text_style(
                text_box.text_frame, item_text,
                font_size=FontSizes.STRATEGY_ITEM,
                font_color=text_color,
                word_wrap=True  # Enable word wrapping
            )
            shapes.append(text_box)
            
            current_y += item_height + item_gap
        
        # Final statement - CSS: font-size: 16px
        final = element_data.get('final', {})
        if final:
            arrow_text = final.get('arrow', '')
            final_text = final.get('text', '')
            
            final_box = slide.shapes.add_textbox(
                x + pad_x, y + h - pad_y - final_height,
                w - pad_x * 2, final_height
            )
            apply_text_style(
                final_box.text_frame, f"{arrow_text} {final_text}",
                font_size=FontSizes.STRATEGY_FINAL,
                font_color=Colors.ACCENT_CORAL,
                bold=True,
                alignment=PP_ALIGN.CENTER
            )
            shapes.append(final_box)
        
        # Queue animations
        for i, shape in enumerate(shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * 40, duration)
    
    def render_ending(self, slide, element_data: dict, slide_data: dict):
        """Render ending slide
        CSS: .ending-icon { font-size: 48px; margin-bottom: 10px; }
        CSS: .ending-message { font-size: 20px; margin-bottom: 16px; }
        CSS: .ending-thanks { font-size: 28px; }
        """
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        layout, height = self.layout.ending_layout()
        self.layout.flow.advance(height)
        
        # Icon - CSS: font-size: 48px
        icon_text = element_data.get('icon', '*')
        icon_rect = layout['icon']
        icon_box = slide.shapes.add_textbox(*icon_rect.to_inches())
        apply_text_style(
            icon_box.text_frame, icon_text,
            font_size=FontSizes.ENDING_ICON,
            font_color=Colors.ACCENT_CORAL,
            alignment=PP_ALIGN.CENTER
        )
        self.anim_builder.queue_entrance(icon_box, preset, step, 0, duration)
        
        # Message - CSS: font-size: 20px; gradient text
        message_text = element_data.get('message', '')
        message_rect = layout['message']
        message_box = slide.shapes.add_textbox(*message_rect.to_inches())
        apply_text_style(
            message_box.text_frame, message_text,
            font_size=FontSizes.ENDING_MESSAGE,
            font_color=Colors.ACCENT_CORAL,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        self.anim_builder.queue_entrance(message_box, 'slideUp', step, 200, duration)
        
        # Thanks - CSS: font-size: 28px; color: var(--text-tertiary)
        thanks_text = element_data.get('thanks', '')
        thanks_rect = layout['thanks']
        thanks_box = slide.shapes.add_textbox(*thanks_rect.to_inches())
        apply_text_style(
            thanks_box.text_frame, thanks_text,
            font_size=FontSizes.ENDING_THANKS,
            font_color=Colors.TEXT_TERTIARY,
            alignment=PP_ALIGN.CENTER
        )
        self.anim_builder.queue_entrance(thanks_box, 'fadeIn', step, 400, duration)
    
    # ===== Image Component =====
    
    def render_image(self, slide, element_data: dict, slide_data: dict):
        """Render image component (single, grid, or side layout)"""
        layout_type = element_data.get('layout', 'single')
        
        if layout_type == 'grid':
            self._render_image_grid(slide, element_data, slide_data)
        elif layout_type == 'side':
            self._render_image_side(slide, element_data, slide_data)
        else:
            self._render_image_single(slide, element_data, slide_data)
    
    def _render_image_single(self, slide, element_data: dict, slide_data: dict):
        """Render single image with optional caption"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        src = element_data.get('src', '')
        caption = element_data.get('caption', '')

        layout, height = self.layout.image_single(has_caption=bool(caption))
        self.layout.flow.advance(height)
        shapes = []
        
        # Add image
        img_rect = layout['image']
        img_shape = self._add_image_or_placeholder(
            slide, src, img_rect, element_data.get('alt', '')
        )
        if img_shape:
            shapes.append(img_shape)
        
        # Add caption if present - .image-caption: 12px, text-tertiary, italic
        if caption and layout['caption']:
            cap_rect = layout['caption']
            cap_box = slide.shapes.add_textbox(*cap_rect.to_inches())
            apply_text_style(
                cap_box.text_frame, caption,
                font_size=FontSizes.IMAGE_CAPTION,
                font_color=Colors.TEXT_TERTIARY,
                italic=True,
                alignment=PP_ALIGN.CENTER
            )
            shapes.append(cap_box)
        
        # Queue animations
        for i, shape in enumerate(shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * 50, duration)
    
    def _render_image_grid(self, slide, element_data: dict, slide_data: dict):
        """Render image grid - new style: images are spread all over the place and titles are covered at the bottom"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        stagger = anim_config.get('stagger', 0) or 100
        
        images = element_data.get('images', [])
        columns = element_data.get('columns', min(len(images), 4))
        
        layouts, height = self.layout.image_grid(len(images), columns)
        self.layout.flow.advance(height)
        
        for i, (img_data, rect_layout) in enumerate(zip(images, layouts)):
            base_delay = i * stagger
            
            # Image - spread over the entire project area, with rounded corners
            img_rect = rect_layout['image']
            img_shape = self._add_image_or_placeholder(
                slide, img_data.get('src', ''), img_rect, img_data.get('alt', ''),
                rounded=True
            )
            if img_shape:
                self.anim_builder.queue_entrance(img_shape, preset, step, 
                                                 base_delay, duration)
            
            # Title overlay - covers the bottom of the image, white text, only rounded corners at the bottom
            title = img_data.get('title', '')
            if title:
                title_rect = rect_layout['title_overlay']
                img_rect_data = rect_layout['image']
                x, y, w, h = title_rect.to_inches()
                
                # Calculate the fillet percentage to make it consistent with the rounded arc of the image
                img_height = img_rect_data.height
                overlay_height = title_rect.height
                overlay_radius = int(8 * img_height / overlay_height) if overlay_height > 0 else 8
                overlay_radius = min(overlay_radius, 50)
                
                # Translucent background strip (first created with a rectangle, then modified to have rounded corners at the bottom)
                overlay_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                # First modify the geometry to have only bottom rounded corners
                self._set_bottom_rounded_corners(overlay_bg, overlay_radius)
                # Set translucent black background
                overlay_bg.fill.solid()
                overlay_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
                overlay_bg.line.fill.background()  # No borders
                # Set transparency (after fill setting)
                self._set_fill_transparency(overlay_bg, 50)
                self.anim_builder.queue_entrance(overlay_bg, preset, step,
                                                 base_delay + 20, duration)
                
                # Title text - white, centered
                title_box = slide.shapes.add_textbox(x, y, w, h)
                apply_text_style(
                    title_box.text_frame, title,
                    font_size=FontSizes.IMAGE_GRID_TITLE,
                    font_color=Colors.WHITE,
                    bold=True,
                    alignment=PP_ALIGN.CENTER
                )
                # Center vertically
                title_box.text_frame.paragraphs[0].space_before = Pt(4)
                self.anim_builder.queue_entrance(title_box, preset, step,
                                                 base_delay + 30, duration)
    
    def _set_fill_transparency(self, shape, transparency_percent: int):
        """Set the transparency of the shape fill (0-100)"""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        # Get spPr element
        spPr = shape._element.spPr
        
        # Find the solidFill element
        solidFill = spPr.find(qn('a:solidFill'))
        
        if solidFill is not None:
            # Get the color element (srgbClr)
            srgbClr = solidFill.find(qn('a:srgbClr'))
            
            if srgbClr is not None:
                # Remove existing alpha setting
                for alpha in list(srgbClr.findall(qn('a:alpha'))):
                    srgbClr.remove(alpha)
                
                # Add transparency attribute (alpha)
                # PowerPoint alpha: 100000 = 100% opaque, 0 = fully transparent
                alpha_val = (100 - transparency_percent) * 1000
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', str(alpha_val))
    
    def _set_rounded_corners(self, shape, radius_percent: int):
        """Set the corner size of the rounded rectangle"""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        # Get spPr (shape properties) element
        spPr = shape._element.spPr
        
        # Find or create a prstGeom element
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            # Find or create avLst
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            else:
                # Clear existing adjustment values
                for child in list(avLst):
                    avLst.remove(child)
            
            # Set fillet size
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius_percent * 1000}')
    
    def _set_bottom_rounded_corners(self, shape, radius_percent: int):
        """Set the shape to have only the bottom two corners rounded (using round2SameRect)"""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        spPr = shape._element.spPr
        
        # Remove existing geometry definition
        for child in list(spPr):
            if child.tag.endswith('}prstGeom') or child.tag.endswith('}custGeom'):
                spPr.remove(child)
        
        # Create a round2SameRect geometry (two corners on the same side are rounded)
        prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'round2SameRect')
        
        # Set fillet size
        # round2SameRect has two adjustment values: adj1 (top) and adj2 (bottom)
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        
        # adj1 = top rounded corners, set to 0 (no rounded corners)
        gd1 = etree.SubElement(avLst, qn('a:gd'))
        gd1.set('name', 'adj1')
        gd1.set('fmla', 'val 0')
        
        # adj2 = bottom rounded corner
        gd2 = etree.SubElement(avLst, qn('a:gd'))
        gd2.set('name', 'adj2')
        gd2.set('fmla', f'val {radius_percent * 1000}')
    
    def _render_image_side(self, slide, element_data: dict, slide_data: dict):
        """Render image with side-by-side content"""
        anim_config = self.resolver.resolve(element_data, slide_data)
        step = element_data.get('step', 1)
        preset = anim_config['preset']
        duration = anim_config['duration']
        
        src = element_data.get('src', '')
        img_position = element_data.get('imgPosition', 'left')
        title = element_data.get('title', '')
        content = element_data.get('content', '')
        
        layout, height = self.layout.image_side(img_position)
        self.layout.flow.advance(height)
        shapes = []
        
        # Image
        img_rect = layout['image']
        img_shape = self._add_image_or_placeholder(slide, src, img_rect, '')
        if img_shape:
            shapes.append(img_shape)
        
        # Content area
        content_rect = layout['content']
        x, y, w, h = content_rect.to_inches()
        
        # inner margin
        pad = Inches(px_to_inch(16))
        title_height = Inches(px_to_inch(28))
        title_gap = Inches(px_to_inch(8))
        
        # Content background
        content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        apply_shape_fill(content_bg, Colors.BG_CARD, Colors.BORDER)
        shapes.append(content_bg)
        
        current_y = y + pad
        
        # Title - .image-side-title: 18px, font-weight: 600
        if title:
            title_box = slide.shapes.add_textbox(
                x + pad, current_y,
                w - pad * 2, title_height
            )
            apply_text_style(
                title_box.text_frame, title,
                font_size=FontSizes.IMAGE_SIDE_TITLE,
                font_color=Colors.TEXT_PRIMARY,
                bold=True
            )
            shapes.append(title_box)
            current_y += title_height + title_gap
        
        # Content text - .image-side-text: 15px, text-secondary
        if content:
            content_box = slide.shapes.add_textbox(
                x + pad, current_y,
                w - pad * 2, h - (current_y - y) - pad
            )
            apply_text_style(
                content_box.text_frame, content,
                font_size=FontSizes.IMAGE_SIDE_TEXT,
                font_color=Colors.TEXT_SECONDARY
            )
            shapes.append(content_box)
        
        # Queue animations
        for i, shape in enumerate(shapes):
            self.anim_builder.queue_entrance(shape, preset, step, i * 50, duration)
    
    def _add_image_or_placeholder(self, slide, src: str, rect, alt: str = '', rounded: bool = True):
        """Add an image to the slide, or a placeholder if image cannot be loaded"""
        x, y, w, h = rect.to_inches()
        
        # Try to load the image
        image_stream = self._fetch_image(src)
        
        if image_stream:
            try:
                pic = slide.shapes.add_picture(image_stream, x, y, w, h)
                
                # Add rounded corners
                if rounded:
                    self._apply_rounded_corners(pic)
                
                return pic
            except Exception as e:
                print(f"Warning: Failed to add image '{src}': {e}")
        
        # Fallback: create a placeholder shape
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        apply_shape_fill(placeholder, Colors.BG_TERTIARY, Colors.BORDER)
        
        # Add placeholder text
        placeholder.text_frame.paragraphs[0].text = alt or "[Image]"
        placeholder.text_frame.paragraphs[0].font.size = Pt(10)
        placeholder.text_frame.paragraphs[0].font.color.rgb = Colors.TEXT_TERTIARY
        placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return placeholder
    
    def _apply_rounded_corners(self, pic, radius_percent: int = 8):
        """Add rounded corners to images
        Set the image's geometry to a rounded rectangle by modifying XML"""
        from pptx.oxml.ns import qn
        from lxml import etree
        
        # Get the spPr (shape properties) element of the image
        spPr = pic._element.spPr
        
        # Create a prstGeom (preset geometry) element and set it to a rounded rectangle
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }
        
        # Remove existing geometry definition (if any)
        for child in list(spPr):
            if child.tag.endswith('}prstGeom') or child.tag.endswith('}custGeom'):
                spPr.remove(child)
        
        # Create a rounded rectangular geometric shape
        prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
        prstGeom.set('prst', 'roundRect')
        
        # Set fillet size (avLst = adjust value list)
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        # Fillet size, 50000 = 50% (maximum fillet), the smaller the value, the smaller the fillet
        gd.set('fmla', f'val {radius_percent * 1000}')
    
    def _fetch_image(self, src: str):
        """Fetch image from URL or local path, return BytesIO stream or None"""
        if not src:
            return None
        
        try:
            # Check if it's a URL
            if src.startswith(('http://', 'https://')):
                # Download from URL
                req = urllib.request.Request(src, headers={
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                })
                with urllib.request.urlopen(req, timeout=10) as response:
                    image_data = response.read()
                    return io.BytesIO(image_data)
            else:
                # Local file - try multiple resolution strategies
                path = Path(src)
                
                # Strategy 1: Absolute path or path relative to current working directory
                if path.is_absolute() and path.exists():
                    print(f"Loaded image: {path}")
                    return io.BytesIO(path.read_bytes())
                
                # Strategy 2: Path relative to base_path (usually the HTML directory)
                if self.base_path:
                    resolved_path = (self.base_path / src).resolve()
                    if resolved_path.exists():
                        print(f"Loaded image: {resolved_path}")
                        return io.BytesIO(resolved_path.read_bytes())
                
                # Strategy 3: Check if path exists as relative to CWD
                if path.exists():
                    print(f"Loaded image: {path.resolve()}")
                    return io.BytesIO(path.read_bytes())
                
                # Not found
                print(f"Warning: Local image not found: {src}")
                if self.base_path:
                    print(f"  Tried: {(self.base_path / src).resolve()}")
                return None
        except urllib.error.URLError as e:
            print(f"Warning: Failed to download image '{src}': {e}")
            return None
        except Exception as e:
            print(f"Warning: Error loading image '{src}': {e}")
            return None
