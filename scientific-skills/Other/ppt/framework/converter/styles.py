# -*- coding: utf-8 -*-
"""Style mappings for converting CSS to PowerPoint styles

Keep in sync with themes/default/variables.css"""

from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# Default CSS variable values from the HTML framework
# Synchronized from themes/default/variables.css
DEFAULT_CSS_VARIABLES = {
    # Background colors - Claude style warm off-white tones
    'bg-primary': '#FAFAF8',           # Warm off-white, Claude’s iconic background
    'bg-secondary': '#F5F3F0',         # Slightly darker off-white, used for the outer background
    'bg-tertiary': '#EFEDE8',          # One level deeper, used for card background contrast
    'bg-card': '#FFFFFF',              # Pure white card background
    
    # Text colors
    'text-primary': '#1A1A1A',         # Near black to ensure readability
    'text-secondary': '#4A4A4A',       # Dark gray, used for text content
    'text-tertiary': '#8B7355',        # Warm brown, used for secondary/secondary information
    'text-muted': '#6A6A6A',           # Medium gray, used to describe text
    'text-light': '#FAFAF8',           # Off-white, for text on dark backgrounds
    
    # Accent colors
    'accent-coral': '#D97757',         # Warm orange, the main color of Claude brand
    'accent-orange': '#E8A87C',        # Light orange, used to highlight key points
    'accent-brown': '#8B7355',         # Warm brown, used for supporting information
    'accent-blue': '#5a8fd4',
    'accent-green': '#5a9a6a',
    'accent-purple': '#8a7ac4',
    'accent-red': '#d45a5a',
    
    # Border colors - Convert to solid colors based on rgba(0,0,0,x)
    'border-color': '#F0F0F0',         # rgba(0, 0, 0, 0.06) is approximately equal to this value on a white background
    'border-strong': '#E6E6E6',        # rgba(0, 0, 0, 0.10) is approximately equal to this value on a white background
    
    # Terminal colors
    'terminal-bg': '#3d3d3d',
    'terminal-header': '#2d2d2d',
    'terminal-text': '#e0e0e0',
    
    # Special
    'dark-bg': '#1a1a1a',
    'white': '#ffffff',
    
    # Fonts
    'font-sans': 'Inter',              # main font
    'font-display': 'Inter',           # Title font
    'font-mono': 'Consolas',           # monospaced font
}


class StyleMapper:
    """
    Maps CSS styles to PowerPoint styles.
    Can be initialized with CSS variables parsed from HTML.
    """
    
    def __init__(self, css_variables: dict = None):
        self.variables = {**DEFAULT_CSS_VARIABLES}
        if css_variables:
            self.variables.update(css_variables)
        
        # Parse colors to RGBColor objects
        self._colors = {}
        for name, value in self.variables.items():
            if value.startswith('#'):
                self._colors[name] = self._parse_color(value)
    
    def _parse_color(self, hex_value: str) -> RGBColor:
        """Parse hex color to RGBColor"""
        hex_value = hex_value.lstrip('#')
        if len(hex_value) == 3:
            hex_value = ''.join([c*2 for c in hex_value])
        return RGBColor(
            int(hex_value[0:2], 16),
            int(hex_value[2:4], 16),
            int(hex_value[4:6], 16)
        )
    
    def color(self, name: str) -> RGBColor:
        """Get an RGBColor by CSS variable name"""
        if name in self._colors:
            return self._colors[name]
        # Try with dashes
        name_dashed = name.replace('_', '-')
        if name_dashed in self._colors:
            return self._colors[name_dashed]
        # Default to black
        return RGBColor(0, 0, 0)
    
    def hex_color(self, name: str) -> str:
        """Get hex color string by CSS variable name"""
        if name in self.variables:
            return self.variables[name].lstrip('#')
        name_dashed = name.replace('_', '-')
        if name_dashed in self.variables:
            return self.variables[name_dashed].lstrip('#')
        return '000000'
    
    def font(self, name: str) -> str:
        """Get font family name"""
        if name == 'mono' or name == 'font-mono':
            return 'Consolas'
        elif name == 'display' or name == 'font-display':
            return 'Inter'
        else:
            # Inter (sans) is used by default
            return 'Inter'


# ===== Pre-defined Style Presets =====
# Synced from themes/default/components.css

class Fonts:
    """Font configurations - synchronized from variables.css"""
    SANS = 'Inter'                     # Main font (--font-sans)
    DISPLAY = 'Inter'                  # Title font (--font-display)
    MONO = 'Consolas'                  # Monospace font (--font-mono)
    # Fallback font (when Inter is not available)
    SANS_FALLBACK = 'Segoe UI'
    CHINESE = 'Microsoft YaHei'        # Chinese fallback


class FontSizes:
    """Font sizes in points - synced from components.css
    
    Conversion formula: Design size 1280px = 10 inches = 720pt
    1px = 720/1280 = 0.5625pt
    
    However, the display effect of PowerPoint is different from that of browsers, and the coefficients need to be fine-tuned.
    Actually use a coefficient of 0.58-0.62 to maintain visual consistency"""
    
    # Hero mode (large) - .slide-title, .slide-subtitle, etc.
    HERO_TITLE = Pt(21)                # .slide-title: 36px * 0.58 ≈ 21pt
    HERO_SUBTITLE = Pt(12)             # .slide-subtitle: 20px * 0.58 ≈ 12pt
    HERO_BADGE = Pt(7)                 # .slide-badge: 12px * 0.58 ≈ 7pt
    HERO_HINT = Pt(8)                  # .click-hint: 13px * 0.58 ≈ 8pt
    HERO_AUTHOR = Pt(8)                # .slide-author: 14px * 0.58 ≈ 8pt
    
    # Content mode (smaller) - scaled header (scale 0.75)
    CONTENT_TITLE = Pt(16)             # 36px * 0.75 * 0.58 ≈ 16pt
    CONTENT_SUBTITLE = Pt(9)           # 20px * 0.75 * 0.58 ≈ 9pt
    CONTENT_BADGE = Pt(5)              # 12px * 0.75 * 0.58 ≈ 5pt
    
    # Component text - individual component styles (using a 0.58 factor)
    CARD_TITLE = Pt(9)                 # .comparison-card h4: 16px → 9pt
    CARD_TEXT = Pt(8)                  # .card-desc: 14px → 8pt
    CARD_ITEM = Pt(8)                  # .comparison-list li: 14px → 8pt
    
    QUOTE_TEXT = Pt(11)                # .quote-text: 18px → 10pt (slightly larger to stay visible)
    QUOTE_AUTHOR = Pt(7)               # .quote-author: 13px → 7pt
    
    TERMINAL_TEXT = Pt(8)              # .terminal-body: 13px → 8pt
    
    TIMELINE_TITLE = Pt(8)             # .timeline-title: 14px → 8pt
    TIMELINE_DESC = Pt(7)              # .timeline-desc: 12px → 7pt
    TIMELINE_BADGE = Pt(6)             # .timeline-badge: 11px → 6pt
    TIMELINE_CASE = Pt(6)              # .timeline-case: 11px → 6pt
    TIMELINE_ICON = Pt(14)             # .timeline-icon: 24px → 14pt
    
    STATS_VALUE = Pt(13)               # .stat-value: 22px → 13pt (keep it eye-catching)
    STATS_LABEL = Pt(6)                # .stat-label: 11px → 6pt
    STATS_ICON = Pt(11)                # .stat-icon: 18px → 11pt
    STATS_TITLE = Pt(8)                # .stats-title: 14px → 8pt
    
    VALUE_CARD_TITLE = Pt(9)           # .value-card h4: 15px → 9pt
    VALUE_CARD_DESC = Pt(7)            # .value-card p: 13px → 7pt
    VALUE_CARD_ICON = Pt(12)           # .value-card-icon: 20px → 12pt
    
    ASSUMPTION_OLD = Pt(8)             # .assumption-old: 14px → 8pt
    ASSUMPTION_NEW = Pt(8)             # .assumption-new: 14px → 8pt
    
    COMPETITION_TITLE = Pt(8)          # .competition-box h4: 14px → 8pt
    COMPETITION_ITEM = Pt(7)           # .competition-item-box: 13px → 7pt
    COMPETITION_CONCLUSION = Pt(9)     # .competition-conclusion: 16px → 9pt
    
    STRATEGY_TITLE = Pt(9)             # .strategy-box h3: 16px → 9pt
    STRATEGY_ITEM = Pt(8)              # .strategy-item: 14px → 8pt
    STRATEGY_FINAL = Pt(9)             # .strategy-final p: 16px → 9pt
    
    ENDING_ICON = Pt(28)               # .ending-icon: 48px → 28pt
    ENDING_MESSAGE = Pt(12)            # .ending-message: 20px → 12pt
    ENDING_THANKS = Pt(16)             # .ending-thanks: 28px → 16pt
    
    # Image component - picture component
    IMAGE_TITLE = Pt(8)                # .image-title: 14px → 8pt
    IMAGE_CAPTION = Pt(7)              # .image-caption: 12px → 7pt
    IMAGE_GRID_TITLE = Pt(8)           # .image-grid-title: 13px → 8pt
    IMAGE_GRID_CAPTION = Pt(6)         # .image-grid-caption: 11px → 6pt
    IMAGE_SIDE_TITLE = Pt(10)          # .image-side-title: 18px → 10pt
    IMAGE_SIDE_TEXT = Pt(9)            # .image-side-text: 15px → 9pt


class Colors:
    """Pre-defined colors as RGBColor objects - synchronized from variables.css"""
    # Backgrounds - Claude style warm off-white tones
    BG_PRIMARY = RGBColor(0xFA, 0xFA, 0xF8)     # #FAFAF8 warm off-white
    BG_SECONDARY = RGBColor(0xF5, 0xF3, 0xF0)  # #F5F3F0 Slightly darker beige
    BG_TERTIARY = RGBColor(0xEF, 0xED, 0xE8)   # #EFEDE8 Go deeper
    BG_CARD = RGBColor(0xFF, 0xFF, 0xFF)        # #FFFFFF pure white
    
    # Text
    TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x1A)   # #1A1A1A nearly black
    TEXT_SECONDARY = RGBColor(0x4A, 0x4A, 0x4A) # #4A4A4A dark gray
    TEXT_TERTIARY = RGBColor(0x8B, 0x73, 0x55)  # #8B7355 Warm Brown ★Important changes
    TEXT_MUTED = RGBColor(0x6A, 0x6A, 0x6A)     # #6A6A6A medium gray
    TEXT_LIGHT = RGBColor(0xFA, 0xFA, 0xF8)     # #FAFAF8 for dark backgrounds
    
    # Accents - Claude brand colors
    ACCENT_CORAL = RGBColor(0xD9, 0x77, 0x57)   # #D97757 Warm Orange ★Important changes
    ACCENT_ORANGE = RGBColor(0xE8, 0xA8, 0x7C)  # #E8A87C Light Orange ★Important changes
    ACCENT_BROWN = RGBColor(0x8B, 0x73, 0x55)   # #8B7355 warm brown
    ACCENT_BLUE = RGBColor(0x5A, 0x8F, 0xD4)    # #5a8fd4
    ACCENT_GREEN = RGBColor(0x5A, 0x9A, 0x6A)   # #5a9a6a
    ACCENT_PURPLE = RGBColor(0x8A, 0x7A, 0xC4)  # #8a7ac4
    ACCENT_RED = RGBColor(0xD4, 0x5A, 0x5A)     # #d45a5a
    
    # Borders - Approximation calculated based on white background
    BORDER = RGBColor(0xF0, 0xF0, 0xF0)         # rgba(0,0,0,0.06) ≈ #F0F0F0
    BORDER_STRONG = RGBColor(0xE6, 0xE6, 0xE6)  # rgba(0,0,0,0.10) ≈ #E6E6E6
    
    # Terminal
    TERMINAL_BG = RGBColor(0x3D, 0x3D, 0x3D)    # #3d3d3d
    TERMINAL_HEADER = RGBColor(0x2D, 0x2D, 0x2D) # #2d2d2d
    TERMINAL_TEXT = RGBColor(0xE0, 0xE0, 0xE0)  # #e0e0e0
    TERMINAL_PROMPT = RGBColor(0xD9, 0x77, 0x57) # #D97757 (accent-coral)
    TERMINAL_HIGHLIGHT = RGBColor(0xD9, 0x77, 0x57) # #D97757
    
    # Special
    DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)        # #1a1a1a
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # #ffffff
    
    # Strategy box specific (gradient approximation)
    STRATEGY_RIGHT_BG = RGBColor(0x1F, 0x3D, 0x22)  # rgba(52, 168, 83, 0.15) approximation
    
    # Opacity simulation - colors blended with dark background (#1A1A1A) at 50% opacity
    # Used for .strategy-item.wrong { opacity: 0.5; }
    WHITE_DIMMED = RGBColor(0x8D, 0x8D, 0x8D)       # white * 0.5 + #1A1A1A * 0.5
    CORAL_DIMMED = RGBColor(0x7A, 0x49, 0x39)       # #D97757 * 0.5 + #1A1A1A * 0.5


# ===== Helper Functions =====

def apply_text_style(text_frame, text: str, 
                     font_size=None, font_name=None, font_color=None,
                     bold=False, italic=False, alignment=None,
                     word_wrap=True, auto_size=None, line_spacing=None):
    """
    Apply text style to a text frame.
    
    Args:
        text_frame: The python-pptx text frame object
        text: The text content
        font_size: Pt() size or None
        font_name: Font family name or None (defaults to Inter)
        font_color: RGBColor or None
        bold: Whether to make text bold
        italic: Whether to make text italic
        alignment: PP_ALIGN value or None
        word_wrap: Whether to enable word wrapping (default True)
        auto_size: Auto size mode - None, 'fit', or 'shrink'
                   'fit': MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE (shrink text to fit)
                   'shrink': MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT (expand shape)
        line_spacing: Line spacing multiplier (e.g., 1.2 for 120%)
    """
    from pptx.enum.text import MSO_AUTO_SIZE
    
    # Set text box properties
    text_frame.word_wrap = word_wrap
    
    # Set autoscaling mode
    if auto_size == 'fit':
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    elif auto_size == 'shrink':
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    elif auto_size is None:
        # Does not automatically adjust by default, but keeps content visible
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    p = text_frame.paragraphs[0]
    p.text = text
    
    if font_size:
        p.font.size = font_size
    
    # The Inter font is used by default, consistent with the HTML version
    p.font.name = font_name if font_name else Fonts.SANS
    
    if font_color:
        p.font.color.rgb = font_color
    if bold:
        p.font.bold = True
    if italic:
        p.font.italic = True
    if alignment:
        p.alignment = alignment
    
    # Set line spacing
    if line_spacing:
        p.line_spacing = line_spacing


def apply_shape_fill(shape, fill_color: RGBColor, 
                     border_color: RGBColor = None, border_width=None):
    """
    Apply fill and border to a shape.
    
    Args:
        shape: The python-pptx shape object
        fill_color: RGBColor for fill
        border_color: RGBColor for border or None
        border_width: Pt() for border width or None
    """
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    if border_color:
        shape.line.color.rgb = border_color
    if border_width:
        shape.line.width = border_width


def get_shape_type(name: str):
    """Get MSO_SHAPE type by name"""
    types = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'rounded': MSO_SHAPE.ROUNDED_RECTANGLE,
    }
    return types.get(name, MSO_SHAPE.ROUNDED_RECTANGLE)
