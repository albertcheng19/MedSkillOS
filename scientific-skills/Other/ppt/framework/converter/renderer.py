# -*- coding: utf-8 -*-
"""
Main slide renderer implementing the Step system and Hero-to-Content mode
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .parsers import ConfigParser, DataParser
from .animation import AnimationResolver, AnimationBuilder
from .layout import LayoutCalculator, px_to_inch, center_x, parse_style
from .styles import Colors, FontSizes, Fonts, apply_text_style, apply_shape_fill
from .components import ComponentRenderers


class SlideRenderer:
    """
    Main renderer that converts slide data to PowerPoint slides.
    Implements:
    - Hero-to-Content mode transition
    - Step-based animations (step, hide-step, activeStep, revealStep)
    """
    
    def __init__(self, config_parser: ConfigParser, data_parser: DataParser, 
                 html_path: str = None, data_path: str = None):
        self.config = config_parser
        self.data = data_parser
        
        # Store base paths for resolving relative image paths
        self.html_base_dir = Path(html_path).parent if html_path else None
        self.data_base_dir = Path(data_path).parent if data_path else None
        
        # Initialize layout calculator
        self.layout = LayoutCalculator()
        
        # Initialize animation system
        self.resolver = AnimationResolver(self.config.component_defaults)
        self.anim_builder = AnimationBuilder(self.resolver)
        
        # Initialize component renderers with base path for images
        self.components = ComponentRenderers(
            self.layout, self.resolver, self.anim_builder,
            base_path=self.html_base_dir
        )
        
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
    
    def render_all(self) -> Presentation:
        """
        Render all slides and return the Presentation object.
        """
        for i, slide_data in enumerate(self.data.slides):
            print(f"Rendering slide {i + 1}...")
            self._render_slide(slide_data)
        
        return self.prs
    
    def _render_slide(self, slide_data: dict):
        """
        Render a single slide with all its elements.
        
        Implements Hero-to-Content mode:
        - Step 0: Show Hero header (large, centered) + clickHint
        - Step 1: Exit Hero header, enter Content header, show first elements
        - Step 2+: Progressive element reveal
        """
        # Create blank slide
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Reset animation builder for new slide
        self.anim_builder.reset()
        
        # Reset flow layout for new slide
        self.layout.reset_flow()
        
        # Add background
        self._add_background(slide)
        
        # Add header (implements Hero-to-Content mode)
        self._add_header(slide, slide_data)
        
        # Group elements by step and render
        elements = slide_data.get('elements', [])
        self._render_elements_by_step(slide, elements, slide_data)
        
        # Apply all animations to slide
        self.anim_builder.apply_to_slide(slide)
    
    def _render_elements_by_step(self, slide, elements: list, slide_data: dict):
        """Render elements step by step and handle layout changes caused by hideStep.
        
        Core logic: When an element of a step is hidden by hideStep,
        The next new element should appear in the position of the hidden element, not stacked below."""
        # Sort by steps
        sorted_elements = sorted(elements, key=lambda e: (e.get('step', 1), elements.index(e)))
        
        # Record the Y position at the beginning of each step
        step_start_y = {}
        # Record the rendering position of each element (including margin-top information)
        element_positions = {}
        
        current_step = 0
        
        for element in sorted_elements:
            step = element.get('step', 1)
            hide_step = element.get('hideStep')
            
            # Parse the style of the element to get margin-top
            style_str = element.get('style', '')
            margins = parse_style(style_str)
            
            # new step begins
            if step != current_step:
                # Check if any elements are hidden in this step
                # If so, reset the Y position to the starting position of the hidden element (including its margin-top)
                reset_y = self._find_hidden_element_y(
                    sorted_elements, element_positions, step
                )
                
                if reset_y is not None:
                    # Reset fluid layout to the position of hidden elements
                    self.layout.flow.current_y = reset_y
                
                # Record the start Y of this step
                step_start_y[step] = self.layout.flow.current_y
                current_step = step
            
            # Record the Y position of the element before rendering (excluding margin-top)
            element_y_before = self.layout.flow.current_y
            # The actual content starting position of the element (including margin-top)
            element_content_y = element_y_before + margins.top
            
            # render elements
            self.components.render(slide, element, slide_data)
            
            # Record the position information of the element (for reference in subsequent steps)
            # start_y records the position containing margin-top so that subsequent elements can start from the same visual position
            element_positions[id(element)] = {
                'step': step,
                'hide_step': hide_step,
                'start_y': element_content_y,  # Actual content starting position (including margin-top)
                'margin_top': margins.top,
                'end_y': self.layout.flow.current_y
            }
    
    def _find_hidden_element_y(self, all_elements: list, positions: dict, current_step: int) -> float:
        """Find the element that is hidden at the current step and return the starting Y position of the earliest hidden element"""
        hidden_y = None
        
        for elem in all_elements:
            elem_id = id(elem)
            if elem_id not in positions:
                continue
            
            pos = positions[elem_id]
            hide_step = pos.get('hide_step')
            
            # If this element is hidden at the current step
            if hide_step and hide_step == current_step:
                # Record the starting position of the earliest hidden element
                if hidden_y is None or pos['start_y'] < hidden_y:
                    hidden_y = pos['start_y']
        
        return hidden_y
    
    def _add_background(self, slide):
        """Add slide background"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        apply_shape_fill(bg, Colors.BG_PRIMARY, None)
        bg.line.fill.background()  # No border
    
    def _add_header(self, slide, slide_data: dict):
        """
        Add slide header with Hero-to-Content transition.
        
        Creates TWO sets of header elements:
        1. Hero version (step 0): Large, centered - exits at step 1
        2. Content version (step 1+): Small, top - enters at step 1
        """
        badge = slide_data.get('badge', {})
        title = slide_data.get('title', {})
        subtitle = slide_data.get('subtitle', '')
        author = slide_data.get('author', '')
        click_hint = slide_data.get('clickHint', '')
        
        # ===== Hero Version (Step 0 - visible on load) =====
        
        # Hero Badge - .slide-badge uses accent-coral
        hero_badge_rect = self.layout.hero_badge()
        hero_badge = slide.shapes.add_textbox(*hero_badge_rect.to_inches())
        badge_text = f"{badge.get('icon', '')}  {badge.get('text', '')}"
        apply_text_style(
            hero_badge.text_frame, badge_text,
            font_size=FontSizes.HERO_BADGE,
            font_color=Colors.ACCENT_CORAL,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        # Hero badge is visible at step 0, exits at step 1
        self.anim_builder.queue_exit(hero_badge, 'fadeOut', 1, 0, 200)
        
        # Hero Title
        hero_title_rect = self.layout.hero_title()
        hero_title = slide.shapes.add_textbox(*hero_title_rect.to_inches())
        title_text = title.get('text', '')
        title_color = Colors.ACCENT_CORAL if title.get('gradient') else Colors.TEXT_PRIMARY
        apply_text_style(
            hero_title.text_frame, title_text,
            font_size=FontSizes.HERO_TITLE,
            font_color=title_color,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        # Hero title is visible at step 0, exits at step 1
        self.anim_builder.queue_exit(hero_title, 'fadeOut', 1, 50, 200)
        
        # Hero Subtitle (if present) - .slide-subtitle uses text-tertiary (warm brown)
        if subtitle:
            hero_sub_rect = self.layout.hero_subtitle()
            hero_sub = slide.shapes.add_textbox(*hero_sub_rect.to_inches())
            apply_text_style(
                hero_sub.text_frame, subtitle,
                font_size=FontSizes.HERO_SUBTITLE,
                font_color=Colors.TEXT_TERTIARY,
                alignment=PP_ALIGN.CENTER
            )
            self.anim_builder.queue_exit(hero_sub, 'fadeOut', 1, 100, 200)
        
        # Hero Author (if present) - displayed below subtitle
        if author:
            hero_author_rect = self.layout.hero_author()
            hero_author = slide.shapes.add_textbox(*hero_author_rect.to_inches())
            apply_text_style(
                hero_author.text_frame, author,
                font_size=FontSizes.HERO_HINT,
                font_color=Colors.TEXT_TERTIARY,
                alignment=PP_ALIGN.CENTER
            )
            self.anim_builder.queue_exit(hero_author, 'fadeOut', 1, 120, 200)
        
        # Click Hint (step 0 only)
        if click_hint:
            hint_rect = self.layout.hero_click_hint()
            hint_box = slide.shapes.add_textbox(*hint_rect.to_inches())
            apply_text_style(
                hint_box.text_frame, f'\u25bc  {click_hint}',
                font_size=FontSizes.HERO_HINT,
                font_color=Colors.TEXT_TERTIARY,
                alignment=PP_ALIGN.CENTER
            )
            # Click hint exits at step 1
            self.anim_builder.queue_exit(hint_box, 'fadeOut', 1, 0, 200)
        
        # ===== Content Version (Step 1+ - enters after click) =====
        
        # Content Badge
        content_badge_rect = self.layout.content_badge()
        content_badge = slide.shapes.add_textbox(*content_badge_rect.to_inches())
        apply_text_style(
            content_badge.text_frame, badge_text,
            font_size=FontSizes.CONTENT_BADGE,
            font_color=Colors.ACCENT_CORAL,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        # Content badge enters at step 1
        self.anim_builder.queue_entrance(content_badge, 'slideDown', 1, 200, 300)
        
        # Content Title
        content_title_rect = self.layout.content_title()
        content_title = slide.shapes.add_textbox(*content_title_rect.to_inches())
        apply_text_style(
            content_title.text_frame, title_text,
            font_size=FontSizes.CONTENT_TITLE,
            font_color=title_color,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        # Content title enters at step 1
        self.anim_builder.queue_entrance(content_title, 'slideDown', 1, 250, 300)
        
        # Content Subtitle (if present) - uses text-tertiary (warm brown)
        if subtitle:
            content_sub_rect = self.layout.content_subtitle()
            content_sub = slide.shapes.add_textbox(*content_sub_rect.to_inches())
            apply_text_style(
                content_sub.text_frame, subtitle,
                font_size=FontSizes.CONTENT_SUBTITLE,
                font_color=Colors.TEXT_TERTIARY,
                alignment=PP_ALIGN.CENTER
            )
            self.anim_builder.queue_entrance(content_sub, 'slideDown', 1, 300, 300)
    
    def save(self, path: str):
        """Save the presentation to a file"""
        self.prs.save(path)
        print(f"Saved presentation to: {path}")


def create_presentation(html_path: str, data_path: str, output_path: str):
    """
    Main function to create a PPTX from HTML/JS sources.
    
    Args:
        html_path: Path to presentation.html
        data_path: Path to slides-data.js
        output_path: Output PPTX file path
    """
    print("Parsing configuration from HTML...")
    config = ConfigParser(html_path)
    
    print("Parsing slide data from JS...")
    data = DataParser(data_path)
    
    print(f"Found {len(data.slides)} slides")
    
    print("Creating presentation...")
    renderer = SlideRenderer(config, data, html_path=html_path, data_path=data_path)
    prs = renderer.render_all()
    
    renderer.save(output_path)
    print("Done!")
    
    return prs
