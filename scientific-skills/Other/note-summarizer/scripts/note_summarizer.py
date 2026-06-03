#!/usr/bin/env python3
"""
Note Summarizer
Summarize notes from Word/PPT/Text/Markdown into a Word outline.
"""

import argparse
import io
import json
from pathlib import Path
from typing import Dict, List, Tuple

from docx import Document
from docx.shared import Inches
from PIL import Image
from pptx import Presentation


DEFAULT_MAX_BULLETS = 6
DEFAULT_MAX_IMAGE_EDGE = 1600


def load_json(path: str) -> Dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def normalize_inputs(inputs: List[str]) -> List[Path]:
    return [Path(p) for p in inputs]


def read_text_file(path: Path) -> List[str]:
    return path.read_text(encoding="utf-8").splitlines()


def extract_docx(path: Path) -> List[Tuple[str, str]]:
    doc = Document(path)
    blocks = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style = p.style.name if p.style else ""
        blocks.append((style, text))
    return blocks


def extract_pptx(path: Path) -> List[Tuple[str, str]]:
    pres = Presentation(path)
    blocks = []
    for slide in pres.slides:
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        if title:
            blocks.append(("Heading 1", title))
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if text and text != title:
                blocks.append(("Body", text))
    return blocks


def extract_images_docx(path: Path) -> List[Image.Image]:
    images = []
    doc = Document(path)
    for rel in doc.part._rels.values():
        if "image" in rel.reltype:
            data = rel.target_part.blob
            img = Image.open(io.BytesIO(data))
            images.append(img.copy())
            img.close()
    return images


def extract_images_pptx(path: Path) -> List[Image.Image]:
    images = []
    pres = Presentation(path)
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                img = Image.open(io.BytesIO(image.blob))
                images.append(img.copy())
                img.close()
    return images


def resize_image(img: Image.Image, max_edge: int) -> Image.Image:
    width, height = img.size
    max_current = max(width, height)
    if max_current <= max_edge:
        return img
    scale = max_edge / max_current
    new_size = (int(width * scale), int(height * scale))
    return img.resize(new_size, Image.LANCZOS)


def to_outline(
    blocks: List[Tuple[str, str]], keep_headings: bool
) -> List[Tuple[str, str]]:
    outline = []
    for style, text in blocks:
        if keep_headings and "Heading" in style:
            outline.append((style, text))
        else:
            outline.append(("Body", text))
    return outline


def summarize_blocks(
    blocks: List[Tuple[str, str]], max_bullets: int
) -> List[Tuple[str, List[str]]]:
    sections = []
    current_title = "Note organization"
    bullets = []

    for style, text in blocks:
        if style.startswith("Heading"):
            if bullets:
                sections.append((current_title, bullets))
            current_title = text
            bullets = []
        else:
            for line in text.split("\n"):
                line = line.strip("- •\t ")
                if not line:
                    continue
                if len(bullets) < max_bullets:
                    bullets.append(line)

    if bullets:
        sections.append((current_title, bullets))

    if not sections:
        sections.append((current_title, []))

    return sections


def add_images(doc: Document, images: List[Image.Image], max_edge: int) -> None:
    for img in images:
        resized = resize_image(img, max_edge)
        stream = io.BytesIO()
        resized.save(stream, format="PNG")
        stream.seek(0)
        width_inch = resized.size[0] / 96
        doc.add_picture(stream, width=Inches(min(width_inch, 6.5)))


def build_document(
    sections: List[Tuple[str, List[str]]], images: List[Image.Image], config: Dict
) -> Document:
    doc = Document()
    doc.add_heading("Study note organization", level=1)

    for title, bullets in sections:
        doc.add_heading(title, level=2)
        if bullets:
            for bullet in bullets:
                doc.add_paragraph(bullet, style="List Bullet")
        else:
            doc.add_paragraph("(No points to extract)")

    if config.get("image_mode") == "copy" and images:
        doc.add_page_break()
        doc.add_heading("Appendix: Pictures", level=2)
        add_images(
            doc, images, config.get("max_image_long_edge", DEFAULT_MAX_IMAGE_EDGE)
        )

    return doc


def collect_interactive() -> Dict:
    inputs = []
    print("Enter input files (blank to finish):")
    while True:
        value = input("Input path: ").strip()
        if not value:
            break
        inputs.append(value)

    output = input("Output .docx path: ").strip() or "summary.docx"
    keep_headings = input("Keep headings? (y/n): ").strip().lower() != "n"
    max_bullets = input("Max bullets per section (default 6): ").strip()
    image_mode = input("Image mode (copy/skip): ").strip().lower() or "copy"
    max_edge = input("Max image long edge (default 1600): ").strip()

    return {
        "inputs": inputs,
        "output": output,
        "keep_headings": keep_headings,
        "max_bullets_per_section": int(max_bullets)
        if max_bullets
        else DEFAULT_MAX_BULLETS,
        "image_mode": image_mode,
        "max_image_long_edge": int(max_edge) if max_edge else DEFAULT_MAX_IMAGE_EDGE,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Note Summarizer")
    parser.add_argument("--json", help="Path to JSON input")
    args = parser.parse_args()

    if args.json:
        config = load_json(args.json)
    else:
        config = collect_interactive()

    inputs = normalize_inputs(config.get("inputs", []))
    keep_headings = config.get("keep_headings", True)
    max_bullets = config.get("max_bullets_per_section", DEFAULT_MAX_BULLETS)
    image_mode = config.get("image_mode", "copy")

    all_blocks: List[Tuple[str, str]] = []
    all_images: List[Image.Image] = []

    for path in inputs:
        if not path.exists():
            continue
        suffix = path.suffix.lower()
        if suffix == ".docx":
            all_blocks.extend(to_outline(extract_docx(path), keep_headings))
            if image_mode == "copy":
                all_images.extend(extract_images_docx(path))
        elif suffix == ".pptx":
            all_blocks.extend(to_outline(extract_pptx(path), keep_headings))
            if image_mode == "copy":
                all_images.extend(extract_images_pptx(path))
        elif suffix in {".txt", ".md"}:
            lines = read_text_file(path)
            blocks = [("Body", line) for line in lines if line.strip()]
            all_blocks.extend(blocks)

    sections = summarize_blocks(all_blocks, max_bullets)
    doc = build_document(sections, all_images, config)

    output_path = Path(config.get("output", "summary.docx"))
    doc.save(output_path)
    print(f"Summary saved to {output_path}")


if __name__ == "__main__":
    main()
